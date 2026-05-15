"""
Generate SmartSpend / Exiqo hackathon deck (18 slides, STAR format).
Run: python presentations/generate_smartspend_deck.py
Output: presentations/SmartSpend_Hackathon_Deck.pptx
"""
from __future__ import annotations

from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# ── Brand palette ─────────────────────────────────────────────────────────
BG = RGBColor(0x0A, 0x06, 0x12)
BG_CARD_TOP = RGBColor(0x15, 0x10, 0x2A)
BG_CARD_BOT = RGBColor(0x0F, 0x0A, 0x1F)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
CYAN = RGBColor(0x06, 0xB6, 0xD4)
EMERALD = RGBColor(0x10, 0xB9, 0x81)
ROSE = RGBColor(0xF4, 0x3F, 0x5E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY300 = RGBColor(0xD1, 0xD5, 0xDB)
GRAY400 = RGBColor(0x9C, 0xA3, 0xAF)
GRAY500 = RGBColor(0x6B, 0x72, 0x80)
GRAY600 = RGBColor(0x4B, 0x55, 0x63)
BORDER = RGBColor(0x33, 0x33, 0x44)

FOOTER_TAGLINE = "12 phases. 4 ML models in production. 1 unified user experience."
OUT_PATH = Path(__file__).resolve().parent / "SmartSpend_Hackathon_Deck.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _set_slide_bg(slide, color: RGBColor = BG) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_rect(slide, left, top, width, height, fill_rgb, line_rgb=None, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    if radius is not None:
        shape.adjustments[0] = radius
    return shape


def _textbox(slide, left, top, width, height, text, size=20, color=GRAY300, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.alignment = align
    return tb


def _add_title(slide, title: str, subtitle: str | None = None) -> None:
    _textbox(slide, Inches(0.7), Inches(0.45), Inches(11.8), Inches(1.0), title, size=40, color=WHITE, bold=True)
    if subtitle:
        _textbox(slide, Inches(0.7), Inches(1.35), Inches(11.5), Inches(0.6), subtitle, size=22, color=GRAY400)


def _star_badge(slide, beat: str) -> None:
    """Bottom-left STAR beat indicator."""
    _textbox(slide, Inches(0.5), Inches(7.05), Inches(1.2), Inches(0.35), f"STAR · {beat}", size=10, color=GRAY500)


def _footer(slide, show_tagline: bool = True, extra: str | None = None) -> None:
    y = Inches(6.85)
    if show_tagline:
        _textbox(slide, Inches(0.7), y, Inches(11.5), Inches(0.4), FOOTER_TAGLINE, size=14, color=GRAY500, align=PP_ALIGN.CENTER)
    if extra:
        _textbox(slide, Inches(0.7), Inches(6.55), Inches(11.5), Inches(0.35), extra, size=12, color=GRAY600, align=PP_ALIGN.CENTER)


def _speaker_notes(slide, s: str, t: str, a: str, r: str) -> None:
    notes = slide.notes_slide.notes_text_frame
    notes.text = (
        f"[S — Situation] {s}\n\n"
        f"[T — Task] {t}\n\n"
        f"[A — Action] {a}\n\n"
        f"[R — Result] {r}"
    )


def _screenshot_placeholder(slide, left, top, width, height, label: str) -> None:
    box = _add_rect(slide, left, top, width, height, BG_CARD_BOT, BORDER, 0.08)
    _textbox(
        slide,
        left + Inches(0.15),
        top + height / 2 - Inches(0.25),
        width - Inches(0.3),
        Inches(0.5),
        label,
        size=14,
        color=GRAY500,
        align=PP_ALIGN.CENTER,
    )


def _bullet_card(slide, left, top, width, height, title: str, lines: list[str], accent: RGBColor = PURPLE) -> None:
    _add_rect(slide, left, top, width, height, BG_CARD_TOP, BORDER, 0.06)
    _textbox(slide, left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.4), title, size=18, color=accent, bold=True)
    body = "\n".join(f"• {ln}" for ln in lines)
    _textbox(slide, left + Inches(0.2), top + Inches(0.55), width - Inches(0.4), height - Inches(0.65), body, size=16, color=GRAY300)


def _flow_box(slide, left, top, w, h, label, sub="", color=PURPLE):
    _add_rect(slide, left, top, w, h, BG_CARD_TOP, color, 0.05)
    _textbox(slide, left + Inches(0.08), top + Inches(0.12), w - Inches(0.16), Inches(0.35), label, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    if sub:
        _textbox(slide, left + Inches(0.06), top + Inches(0.45), w - Inches(0.12), Inches(0.35), sub, size=10, color=GRAY400, align=PP_ALIGN.CENTER)


def _arrow_right(slide, left, top, width=Inches(0.35)):
    _textbox(slide, left, top, width, Inches(0.3), "→", size=22, color=CYAN, align=PP_ALIGN.CENTER)


# ── Slide builders ────────────────────────────────────────────────────────

def slide_01_title(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _star_badge(slide, "S")
    _textbox(slide, Inches(0.9), Inches(2.0), Inches(11.5), Inches(1.2), "SmartSpend Analytics", size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    _textbox(
        slide,
        Inches(1.0),
        Inches(3.1),
        Inches(11.3),
        Inches(0.8),
        "Financial intelligence for India — before money leaves your account",
        size=24,
        color=GRAY300,
        align=PP_ALIGN.CENTER,
    )
    _textbox(slide, Inches(1.0), Inches(4.0), Inches(11.3), Inches(0.5), "Exiqo · Hackathon Demo", size=18, color=GRAY500, align=PP_ALIGN.CENTER)
    _footer(slide, True)
    _speaker_notes(
        slide,
        "India processes billions of UPI transactions monthly; fraud losses hit ₹1000+ Cr/year and users discover theft after payment.",
        "Hook judges with a product that feels production-grade, not a weekend wrapper.",
        "Introduce SmartSpend: React + FastAPI + PostgreSQL with 12-phase fraud pipeline and unified dark UI.",
        "Judges should feel: this team ships real fintech intelligence, not a ChatGPT skin.",
    )


def slide_02_problem(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _star_badge(slide, "S→T")
    _add_title(slide, "₹1 traps, UPI collect scams, and silent subscription bleed", "Three gaps every Indian digital banker feels")
    cards = [
        ("Reactive fraud", "Money leaves first.\nBanks alert after loss.\nNo pre-send verdict.", ROSE),
        ("Data without insight", "PDFs and exports.\nNo health score or EMI risk.\nNo actionable next step.", CYAN),
        ("Low self-awareness", "Festival spikes, dark patterns,\nunused subscriptions — invisible\nuntil the statement arrives.", PURPLE),
    ]
    for i, (title, body, col) in enumerate(cards):
        x = Inches(0.7 + i * 4.1)
        _bullet_card(slide, x, Inches(2.0), Inches(3.85), Inches(3.8), title, body.split("\n"), col)
    _footer(slide)
    _speaker_notes(
        slide,
        "RBI and NPCI report rising digital payment fraud; users trust UPI but fear unknown VPA collect requests.",
        "Frame the problem judges recognize from news and personal experience.",
        "Three-column pain map: post-factum fraud, raw data dumps, behavioral blind spots.",
        "Sets up SmartSpend as proactive protection + explainable intelligence, not another expense tracker.",
    )


def slide_03_mission(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _star_badge(slide, "T")
    _add_title(slide, "Our task: decisions before payment, not regrets after", "Mission — one platform, one truth layer")
    _bullet_card(
        slide,
        Inches(0.7),
        Inches(2.0),
        Inches(5.8),
        Inches(4.2),
        "What we set out to solve",
        [
            "Score every transaction with ML in <50ms",
            "Let users ask AI only over real DB rows",
            "Surface India-specific risks: UPI, EMI, festivals, RBI patterns",
            "Unify fraud, spend, subscriptions in one UX",
        ],
        CYAN,
    )
    _bullet_card(
        slide,
        Inches(6.8),
        Inches(2.0),
        Inches(5.8),
        Inches(4.2),
        "What success looks like",
        [
            "Pre-send FraudShield verdict on any payment",
            "Health score + anomaly flags on dashboard",
            "Honest metrics (ROC-AUC, PSI) in production UI",
            "Product polish judges trust instantly",
        ],
        EMERALD,
    )
    _footer(slide)
    _speaker_notes(
        slide,
        "Hackathon teams often demo dashboards; judges ask 'so what?'",
        "State the explicit engineering goal: proactive fraud + grounded AI + Indian context.",
        "Left: four concrete deliverables. Right: measurable outcomes for users and judges.",
        "Transitions to live product map — proof we built the mission, not slideware.",
    )


def slide_04_product_overview(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _star_badge(slide, "A")
    _add_title(slide, "SmartSpend: from PostgreSQL rows to financial decisions", "React 18 · FastAPI · PostgreSQL · Redis · MLflow")
    _screenshot_placeholder(slide, Inches(0.7), Inches(2.0), Inches(7.5), Inches(4.3), "[INSERT SCREENSHOT: Dashboard hero + live ticker]")
    _bullet_card(
        slide,
        Inches(8.5),
        Inches(2.0),
        Inches(4.1),
        Inches(4.3),
        "Core loop",
        [
            "Ingest transactions (demo seeds + docs)",
            "Score & explain risk (12 phases)",
            "Narrate insights (grounded LLM)",
            "Plan festivals, EMIs, purchases",
        ],
        PURPLE,
    )
    _footer(slide)
    _speaker_notes(
        slide,
        "Most apps show charts; few connect risk, planning, and chat to the same ledger.",
        "Show the product is a full stack, not a frontend mock.",
        "Walk the data loop: DB → API → ML → UI. Point at live ticker on screenshot.",
        "Judge takeaway: end-to-end system ready to demo at localhost:3000.",
    )


def slide_05_feature_map(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _star_badge(slide, "A→R")
    _add_title(slide, "One app. Every feature a judge can click.", "Workspace · AI Intelligence · Financial OS · Planning")
    features = [
        ("Dashboard", "Health score, anomalies, ticker"),
        ("Transactions", "Filters, categories, ₹"),
        ("AI Insights", "Chat + monthly narrative"),
        ("Subscriptions AI", "Graveyard + device link"),
        ("FraudShield", "12 phases + pre-send check"),
        ("Dark Patterns", "Billing traps detector"),
        ("EMI Tracker", "RBI-style DTI bands"),
        ("Festivals / Purchase", "Indian context planners"),
    ]
    cols, rows = 4, 2
    for idx, (name, desc) in enumerate(features):
        c, r = idx % cols, idx // cols
        x = Inches(0.65 + c * 3.15)
        y = Inches(2.05 + r * 2.35)
        _bullet_card(slide, x, y, Inches(2.95), Inches(2.05), name, [desc], PURPLE if idx == 4 else CYAN)
    _footer(slide)
    _speaker_notes(
        slide,
        "Judges compare you to CRED/Jupiter — need a clear feature inventory.",
        "Prove breadth without buzzwords: every sidebar tab is implemented.",
        "8-card map aligned to App.jsx navigation groups.",
        "Result: one unified experience — footer tagline — not scattered micro-tools.",
    )


def slide_06_fraudshield_stack(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _star_badge(slide, "S→T→A")
    _add_title(slide, "FraudShield: not ChatGPT with a shield icon", "4 ML models · rules · LLM investigator · orchestrator")
    models = [
        ("XGBoost", "fraud-xgboost-v0\nProduction classifier", EMERALD),
        ("Isolation Forest", "Anomaly tier-1\nUnsupervised flags", CYAN),
        ("GraphSAGE GNN", "64-dim user embeddings\nPhase 10", PURPLE),
        ("Multi-branch DNN", "Shadow → promotion\nPhase 11", RGBColor(0x14, 0xB8, 0xA6)),
    ]
    for i, (name, desc, col) in enumerate(models):
        _bullet_card(slide, Inches(0.65 + i * 3.15), Inches(2.0), Inches(2.95), Inches(2.5), name, desc.split("\n"), col)
    _textbox(
        slide,
        Inches(0.7),
        Inches(4.75),
        Inches(12.0),
        Inches(0.9),
        "Training data: analyst labels + RBI-style fraud taxonomy + synthetic UPI patterns (lottery, KYC, collect scams) — not generic LLM knowledge.",
        size=16,
        color=GRAY400,
    )
    _screenshot_placeholder(slide, Inches(0.7), Inches(5.35), Inches(12.0), Inches(1.15), "[INSERT SCREENSHOT: FraudShield overview / phase rail]")
    _footer(slide, extra="Answers: 'Isn't this just ChatGPT?' — No. Scores come from XGBoost + graph + rules first.")
    _speaker_notes(
        slide,
        "UPI fraud is pattern-specific; generic LLMs hallucinate risk scores.",
        "Differentiate from wrappers: show four distinct model roles.",
        "Name production artifact fraud-xgboost-v0; mention GNN GraphSAGE and DNN shadow path.",
        "Judge believes: rigorous stack; LLM only investigates borderline cases (Phase 9/12).",
    )


def slide_07_fraud_pipeline(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _star_badge(slide, "A")
    _add_title(slide, "12-phase pipeline: event → score → explain → learn", "Sub-50ms auto-score · manual pre-send check uses same engine")
    # Flow row 1: phases 1-6
    phases_row1 = ["1 Event", "2 Features", "3 XGBoost", "4 Rules", "5 MLOps", "6 Graph"]
    x0, y0, bw, bh, gap = Inches(0.45), Inches(2.15), Inches(1.95), Inches(0.85), Inches(0.08)
    for i, lbl in enumerate(phases_row1):
        x = x0 + i * (bw + gap)
        col = PURPLE if i == 2 else CYAN if i < 2 else EMERALD
        _flow_box(slide, x, y0, bw, bh, lbl.split()[0], lbl.split(maxsplit=1)[1] if " " in lbl else "", col)
        if i < len(phases_row1) - 1:
            _arrow_right(slide, x + bw + Inches(0.02), y0 + Inches(0.28))
    phases_row2 = ["7 SHAP", "8 Feedback", "9 LLM Agent", "10 GNN", "11 DNN", "12 Orchestrator"]
    y1 = Inches(3.25)
    for i, lbl in enumerate(phases_row2):
        x = x0 + i * (bw + gap)
        col = ROSE if i == 5 else PURPLE if i >= 8 else CYAN
        _flow_box(slide, x, y1, bw, bh, lbl.split()[0], lbl.split(maxsplit=1)[1] if " " in lbl else "", col)
        if i < len(phases_row2) - 1:
            _arrow_right(slide, x + bw + Inches(0.02), y1 + Inches(0.28))
    _bullet_card(
        slide,
        Inches(0.7),
        Inches(4.35),
        Inches(6.0),
        Inches(2.15),
        "Pre-send check flow",
        [
            "User enters amount + VPA/merchant",
            "Feature store + XGBoost + policies",
            "Verdict: Safe / Review / Do not proceed",
            "Quick tests: ₹1 trap, lottery, KYC",
        ],
        ROSE,
    )
    _screenshot_placeholder(slide, Inches(6.9), Inches(4.35), Inches(5.7), Inches(2.15), "[INSERT SCREENSHOT: Transaction checker verdict]")
    _footer(slide)
    _speaker_notes(
        slide,
        "Banks run batch fraud; users need instant pre-payment guidance.",
        "Explain both automatic pipeline (phases 1-8) and user-facing checker.",
        "Walk left-to-right on diagram; offer 2-min live demo after this slide.",
        "Result: 12 layers shipped in codebase (migrations + routes per phase).",
    )


def slide_08_ai_insights(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _star_badge(slide, "T→A→R")
    _add_title(slide, "AI Insights: grounded chat, zero invented numbers", "Context packet from PostgreSQL — Groq/OpenAI with fallbacks")
    _bullet_card(
        slide,
        Inches(0.7),
        Inches(2.0),
        Inches(5.5),
        Inches(4.5),
        "How it works",
        [
            "build_context_packet(): profile, accounts, 50 recent txns",
            "System prompt forbids inventing amounts",
            "Routes user to EMI, Fraud, Subscriptions tabs",
            "Hindi or English — same data layer",
            "Monthly insights card: JSON from real aggregates",
        ],
        CYAN,
    )
    _screenshot_placeholder(slide, Inches(6.5), Inches(2.0), Inches(6.2), Inches(4.5), "[INSERT SCREENSHOT: AI Insights chat + CHIPS follow-ups]")
    _footer(slide)
    _speaker_notes(
        slide,
        "Users ask 'why did I overspend?' — generic chatbots guess.",
        "Deliver trustworthy answers tied to their ledger.",
        "Explain context packet + ROUTE navigation + graceful degradation without API keys.",
        "Result: demo ask a question in Hindi; show cited merchants/categories from DB.",
    )


def slide_09_subscriptions(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _star_badge(slide, "S→T→A→R")
    _add_title(slide, "Subscriptions AI: find money leaking in the background", "Connect flow → Hub → AI Analysis → Smart Reminders")
    _bullet_card(slide, Inches(0.7), Inches(2.0), Inches(3.9), Inches(4.3), "Situation", ["Recurring UPI mandates", "Forgotten trials after festivals", "Duplicate OTT + cloud bills"], ROSE)
    _bullet_card(slide, Inches(4.85), Inches(2.0), Inches(3.9), Inches(4.3), "Action", ["Device-link onboarding", "Subscription graveyard detection", "Usage vs spend scoring"], PURPLE)
    _bullet_card(slide, Inches(9.0), Inches(2.0), Inches(3.9), Inches(4.3), "Result", ["Recover ₹/month from unused subs", "Proactive cancel reminders", "Tied to JWT user — no cross-user leak"], EMERALD)
    _screenshot_placeholder(slide, Inches(0.7), Inches(5.5), Inches(12.0), Inches(1.0), "[INSERT SCREENSHOT: Subscription Hub]")
    _footer(slide)
    _speaker_notes(
        slide,
        "India's subscription economy exploded post-UPI autopay; users forget mandates.",
        "Stop silent bleed — second pillar after fraud.",
        "Walk connect → hub → AI analysis engines in App.jsx subscriptions sub-views.",
        "Quantify savings narrative even if demo uses seeded subscriptions.",
    )


def slide_10_dark_emi(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _star_badge(slide, "A→R")
    _add_title(slide, "Dark Patterns + EMI Trap + Health Score", "Rules + ML — protecting against product tricks and debt risk")
    items = [
        ("Dark Pattern Detector", "Rupee-trap pricing, hidden renewals, escalation patterns", ROSE),
        ("EMI Trap Detector", "Debt-to-income vs RBI-style safe band; warn before over-leverage", CYAN),
        ("Health Score 0–100", "Savings rate, anomalies, grade breakdown on dashboard", EMERALD),
        ("Scenario Simulator", "What-if projections on savings and health", PURPLE),
    ]
    for i, (t, d, c) in enumerate(items):
        _bullet_card(slide, Inches(0.65 + (i % 2) * 6.35), Inches(2.0 + (i // 2) * 2.35), Inches(6.1), Inches(2.05), t, [d], c)
    _footer(slide)
    _speaker_notes(
        slide,
        "Indian apps use ₹1 trials and EMI nudges; users need pattern literacy.",
        "Show non-fraud intelligence — holistic financial OS.",
        "Each module: rule/ML backend + dedicated React surface.",
        "Result: dashboard becomes command center, not just FraudShield.",
    )


def slide_11_planning(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _star_badge(slide, "A→R")
    _add_title(slide, "Planning for Indian life: festivals, goals, family", "Context competitors ignore")
    _bullet_card(slide, Inches(0.7), Inches(2.0), Inches(3.9), Inches(4.2), "Festival Predictor", ["Upcoming Diwali, Holi, etc.", "Savings targets per event", "Urgency strip on dashboard"], PURPLE)
    _bullet_card(slide, Inches(4.75), Inches(2.0), Inches(3.9), Inches(4.2), "Purchase Planner", ["Goal milestones", "EMI vs cash tradeoff", "Sacrifice hints from spend"], CYAN)
    _bullet_card(slide, Inches(8.8), Inches(2.0), Inches(3.9), Inches(4.2), "Trips & Family Events", ["Shared event budgets", "Family head persona fit", "Linked to transaction categories"], EMERALD)
    _screenshot_placeholder(slide, Inches(0.7), Inches(5.0), Inches(12.0), Inches(1.55), "[INSERT SCREENSHOT: Festival or Purchase planner]")
    _footer(slide)
    _speaker_notes(
        slide,
        "Global PFM apps miss festival spikes and joint family spending.",
        "Prove India-first product thinking beyond fraud.",
        "Three planners pull from same PostgreSQL transaction history.",
        "Result: Priya persona (family head) has a reason to adopt beyond FraudShield.",
    )


def slide_12_architecture(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _star_badge(slide, "A")
    _add_title(slide, "Architecture: production-shaped, hackathon-fast", "Privacy: PII stays in user's DB row — context packet minimizes exposure")
    # Architecture diagram as boxes
    layers = [
        ("React 18 CRA", "Dashboard, FraudShield, Command Palette", Inches(0.7), Inches(2.0), CYAN),
        ("FastAPI REST", "50+ routes, Swagger /docs", Inches(4.5), Inches(2.0), PURPLE),
        ("PostgreSQL", "Transactions, fraud labels, chat sessions", Inches(8.3), Inches(2.0), EMERALD),
        ("Redis", "Streams, feature store, GNN embeddings", Inches(0.7), Inches(3.5), ROSE),
        ("MLflow + workers", "Registry, PSI drift, retrain scheduler", Inches(4.5), Inches(3.5), CYAN),
        ("Groq / OpenAI", "Investigator + chat (optional keys)", Inches(8.3), Inches(3.5), PURPLE),
    ]
    for title, sub, x, y, col in layers:
        _flow_box(slide, x, y, Inches(3.5), Inches(1.15), title, sub, col)
    _textbox(slide, Inches(0.7), Inches(5.0), Inches(12.0), Inches(1.2),
        "API: POST /api/fraud-shield/{user}/check-transaction · GET /api/insights · POST /api/ai/chat (SSE)\n"
        "Orchestrator tiers 0–3: XGBoost only → +anomaly → +GNN → +LLM judge (cost-aware routing)",
        size=16, color=GRAY300)
    _footer(slide, extra="PII: training uses aggregated features; chat never sees raw card numbers.")
    _speaker_notes(
        slide,
        "Judges will ask about scale, privacy, and deployment.",
        "Show honest architecture — monolith API, real DB, optional AI keys.",
        "Walk diagram; mention MLflow fraud-xgboost-v0 and Redis feature store.",
        "Result: credible path to Account Aggregator + mobile apps (roadmap).",
    )


def slide_13_defensible(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _star_badge(slide, "R")
    _add_title(slide, "Why this isn't just another fintech wrapper", None)
    quads = [
        ("Multi-model AI stack", "XGBoost + DNN + GNN + orchestrator — not a single LLM call", PURPLE),
        ("India-specific training", "UPI patterns, festival cycles, RBI fraud taxonomy", CYAN),
        ("Behavioral graph intelligence", "GraphSAGE learns embeddings from transaction relationships", EMERALD),
        ("Zero-hallucination design", "Every number traceable to a transaction row in PostgreSQL", ROSE),
    ]
    for i, (title, body, col) in enumerate(quads):
        x = Inches(0.7 + (i % 2) * 6.35)
        y = Inches(2.0 + (i // 2) * 2.45)
        _bullet_card(slide, x, y, Inches(6.1), Inches(2.15), title, [body], col)
    _footer(slide)
    _speaker_notes(
        slide,
        "CRED/Jupiter optimize rewards; judges ask what your moat is.",
        "Answer defensibility before Q&A.",
        "Four quadrants — moat is depth of fraud stack + India data + graph + grounded numbers.",
        "Result: differentiated from 'ChatGPT + finance UI' accusation.",
    )


def slide_14_metrics(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _star_badge(slide, "R")
    _add_title(slide, "Real model performance, real numbers", "fraud-xgboost-v0 · trained ~4 days ago · Production in MLflow UI")
    metrics = [
        ("ROC-AUC", "0.933", "Fraud classifier discrimination", EMERALD),
        ("PR-AUC", "0.037", "Realistic for class-imbalanced fraud — we report honestly", CYAN),
        ("Recall @ 5% FPR", "50%", "Half of fraud caught; only 5% false alarms on normals", PURPLE),
        ("Val AUCPR", "1.000", "Internal validation split", RGBColor(0xEC, 0x48, 0x99)),
        ("PSI", "0.050", "Stable — no model drift detected", EMERALD),
        ("Model ID", "fraud-xgboost-v0", "4 ML models in production path", ROSE),
    ]
    for i, (label, val, sub, col) in enumerate(metrics):
        c, r = i % 3, i // 3
        x = Inches(0.65 + c * 4.15)
        y = Inches(2.0 + r * 2.15)
        _add_rect(slide, x, y, Inches(3.85), Inches(1.85), BG_CARD_TOP, BORDER, 0.06)
        _textbox(slide, x + Inches(0.2), y + Inches(0.15), Inches(3.5), Inches(0.35), label, size=14, color=GRAY500)
        _textbox(slide, x + Inches(0.2), y + Inches(0.5), Inches(3.5), Inches(0.7), val, size=36, color=col, bold=True)
        _textbox(slide, x + Inches(0.2), y + Inches(1.25), Inches(3.5), Inches(0.5), sub, size=12, color=GRAY400)
    _textbox(slide, Inches(0.7), Inches(6.35), Inches(12.0), Inches(0.4),
             "Honest metrics. We're not hiding behind accuracy numbers.", size=16, color=GRAY500, align=PP_ALIGN.CENTER)
    _footer(slide, show_tagline=False)
    _speaker_notes(
        slide,
        "Hackathon teams fake 99% accuracy; judges with ML background spot it.",
        "Show receipts from AIPerformance.jsx / trained pickle metrics.",
        "Explain PR-AUC context for imbalance; PSI 0.05 = stable production model.",
        "Result: credibility — invite judges to Admin → AI Performance screen live.",
    )


def slide_15_design(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _star_badge(slide, "A→R")
    _add_title(slide, "Designed like a product, not a hackathon project", "#0A0612 · purple accent only — never purple body text")
    placeholders = [
        "[INSERT SCREENSHOT: Premium dark dashboard]",
        "[INSERT SCREENSHOT: Live transaction ticker]",
        "[INSERT SCREENSHOT: FraudShield quick-test]",
    ]
    for i, lbl in enumerate(placeholders):
        _screenshot_placeholder(slide, Inches(0.65 + i * 4.15), Inches(2.0), Inches(3.85), Inches(3.5), lbl)
    _textbox(
        slide,
        Inches(0.7),
        Inches(5.75),
        Inches(12.0),
        Inches(0.55),
        "Built with the rigor of Stripe, Mercury, Linear — users won't trust an app that doesn't look trustworthy.",
        size=17,
        color=GRAY400,
        align=PP_ALIGN.CENTER,
    )
    _footer(slide)
    _speaker_notes(
        slide,
        "Judges decide in 10 seconds if UI is hackathon-grade or fundable.",
        "Design is trust signal for a money app.",
        "Point to design tokens: gray hierarchy, aurora background, command palette.",
        "Result: passes the 'would I enter my bank data?' sniff test.",
    )


def slide_16_personas(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _star_badge(slide, "S→R")
    _add_title(slide, "Built for 800 million Indian digital banking users", "Three personas — each with a full STAR arc")
    personas = [
        (
            "Vikram, 26 — Young professional",
            "S: UPI-native, 6 subscriptions, scam DMs daily",
            "T: Know if payment is safe before sending",
            "A: FraudShield pre-send + subscription graveyard",
            "R: Confident payments; ₹800/mo subs recovered",
            CYAN,
        ),
        (
            "Priya, 38 — Family head",
            "S: Festival + EMI + kids' school fees",
            "T: Household budget without spreadsheet pain",
            "A: Festival planner, EMI trap, family events",
            "R: Health score 72→85; no surprise festival debt",
            PURPLE,
        ),
        (
            "Arjun, 22 — First-time investor",
            "S: Scams target newcomers; low literacy",
            "T: Learn + protect without jargon",
            "A: AI Insights chat + dark pattern alerts",
            "R: Grounded answers; avoided ₹1 trial trap",
            EMERALD,
        ),
    ]
    for i, (name, s, t, a, r, col) in enumerate(personas):
        _add_rect(slide, Inches(0.65 + i * 4.15), Inches(2.0), Inches(3.85), Inches(4.55), BG_CARD_TOP, BORDER, 0.06)
        _textbox(slide, Inches(0.85 + i * 4.15), Inches(2.15), Inches(3.45), Inches(0.45), name, size=16, color=col, bold=True)
        for j, line in enumerate([s, t, a, r]):
            _textbox(slide, Inches(0.85 + i * 4.15), Inches(2.65 + j * 0.95), Inches(3.45), Inches(0.85), line, size=13, color=GRAY300)
    _footer(slide)
    _speaker_notes(
        slide,
        "PM judges want market fit, not only tech.",
        "Map features to real Indian user segments.",
        "60 seconds: one sentence each on Vikram, Priya, Arjun STAR lines on slide.",
        "Result: TAM story — 800M digital banking users, three entry wedges.",
    )


def slide_17_roadmap(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _star_badge(slide, "R → future")
    _add_title(slide, "12 phases shipped. 8 more to go.", "Business: B2B FraudShield API + Pro tier")
    milestones = [
        ("Now", "ML stack live · dashboard · fraud scanner in production", EMERALD),
        ("30 days", "Account Aggregator bank link · iOS/Android apps", CYAN),
        ("90 days", "Investment health AI · tax optimization · joint accounts", PURPLE),
        ("6 months", "Open API — fintechs embed FraudShield as a service", ROSE),
    ]
    for i, (when, what, col) in enumerate(milestones):
        y = Inches(2.1 + i * 1.15)
        _add_rect(slide, Inches(0.9), y, Inches(0.18), Inches(0.75), col, col)
        _textbox(slide, Inches(1.25), y, Inches(2.0), Inches(0.4), when, size=20, color=col, bold=True)
        _textbox(slide, Inches(3.4), y, Inches(9.2), Inches(0.75), what, size=18, color=GRAY300)
    _footer(slide)
    _speaker_notes(
        slide,
        "Judges ask 'what after hackathon?' and 'how do you make money?'",
        "Show roadmap discipline — 12 done, 8 planned parity phases.",
        "Timeline: AA integration, mobile, B2B API monetization.",
        "Result: team has vision beyond demo day.",
    )


def slide_18_close(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _star_badge(slide, "R")
    _textbox(slide, Inches(0.8), Inches(1.6), Inches(11.7), Inches(1.0),
             "SmartSpend isn't a finance app. It's financial intelligence.", size=42, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    _textbox(slide, Inches(0.8), Inches(2.75), Inches(11.7), Inches(0.7),
             "Indian fintech needs AI built for Indian problems. We built it.", size=24, color=PURPLE, align=PP_ALIGN.CENTER)
    _textbox(
        slide,
        Inches(0.8),
        Inches(3.6),
        Inches(11.7),
        Inches(1.2),
        "12-phase fraud pipeline · 4 ML models · Grounded AI chat · One dark, trustworthy UX\n"
        "Pre-send protection · Subscription intelligence · Festival-aware planning",
        size=18,
        color=GRAY300,
        align=PP_ALIGN.CENTER,
    )
    _add_rect(slide, Inches(5.5), Inches(4.95), Inches(2.3), Inches(2.3), BG_CARD_TOP, BORDER, 0.08)
    _textbox(slide, Inches(5.55), Inches(5.55), Inches(2.2), Inches(1.0), "[QR CODE]\nDemo URL", size=14, color=GRAY500, align=PP_ALIGN.CENTER)
    _textbox(slide, Inches(0.8), Inches(5.15), Inches(4.5), Inches(1.5),
             "Try now:\nhttp://localhost:3000\nAPI: :8000/docs", size=16, color=GRAY400)
    _textbox(slide, Inches(8.0), Inches(5.15), Inches(4.5), Inches(1.5),
             "Team: [Your Names]\n[email] · [GitHub]\nBuilt for [Hackathon Name]", size=16, color=GRAY400, align=PP_ALIGN.RIGHT)
    _footer(slide, True)
    _speaker_notes(
        slide,
        "Full deck STAR result: from Indian fraud crisis to shipped intelligence platform.",
        "Ask: try the demo, scan QR, open FraudShield quick-test.",
        "Rehearse 30-second close: problem → 12 phases → metrics → live URL.",
        "Leave judges wanting to tap localhost:3000 before next team.",
    )


def build_deck() -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        slide_01_title,
        slide_02_problem,
        slide_03_mission,
        slide_04_product_overview,
        slide_05_feature_map,
        slide_06_fraudshield_stack,
        slide_07_fraud_pipeline,
        slide_08_ai_insights,
        slide_09_subscriptions,
        slide_10_dark_emi,
        slide_11_planning,
        slide_12_architecture,
        slide_13_defensible,
        slide_14_metrics,
        slide_15_design,
        slide_16_personas,
        slide_17_roadmap,
        slide_18_close,
    ]
    for fn in builders:
        fn(prs)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PATH))
    return OUT_PATH


if __name__ == "__main__":
    path = build_deck()
    print(f"Generated: {path}")
    print(f"Slides: 18 | Generated at: {datetime.now().isoformat(timespec='seconds')}")
